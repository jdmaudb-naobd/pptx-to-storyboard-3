"""
Simple PowerPoint content extractor
No complex dependencies!
"""

from pptx import Presentation
import json
from pathlib import Path
from src.utils import sanitize_text

class SimpleExtractor:
    def __init__(self, pptx_path):
        self.pptx_path = Path(pptx_path)
        self.presentation = Presentation(str(pptx_path))
        
    def prompt_for_abbreviation_definitions(self, undefined_abbreviations, json_path, content):
        """Prompt user to input definitions for undefined abbreviations, show context, and update JSON library."""
        if not undefined_abbreviations:
            print("No undefined abbreviations found.")
            return

        print("\nUndefined abbreviations detected:")
        new_definitions = {}
        ignored_abbreviations = set()

        for abbr in undefined_abbreviations:
            # Find context for the abbreviation
            context_found = False
            for slide in content["slides"]:
                for text_item in slide["texts"]:
                    words = text_item["text"].split()
                    if abbr in words:
                        idx = words.index(abbr)
                        start_idx = max(0, idx - 10)
                        end_idx = min(len(words), idx + 10)
                        context = " ".join(words[start_idx:end_idx])
                        print(f"Context for '{abbr}': {context}")
                        context_found = True
                        break
                if context_found:
                    break

            action = input(f"Enter definition for '{abbr}' or type 'ignore' to skip: ")
            if action.lower() == 'ignore':
                ignored_abbreviations.add(abbr)
            elif action:
                new_definitions[abbr] = action

        # Update the JSON library
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                existing_definitions = json.load(f)
        except FileNotFoundError:
            existing_definitions = {}

        existing_definitions.update(new_definitions)

        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(existing_definitions, f, indent=2, ensure_ascii=False)

        print("Abbreviation definitions updated successfully.")

        # Remove ignored abbreviations from undefined_abbreviations
        undefined_abbreviations.difference_update(ignored_abbreviations)

    def extract_all_content(self):
        """Extract all content from PowerPoint and handle undefined abbreviations."""
        content = {
            "filename": self.pptx_path.name,
            "slide_count": len(self.presentation.slides),
            "slides": []
        }

        undefined_abbreviations = set()

        for idx, slide in enumerate(self.presentation.slides):
            slide_content = {
                "slide_number": idx + 1,
                "texts": [],
                "shapes": []
            }

            # Extract text from each shape
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    extracted_text = shape.text
                    # Sanitize the extracted text
                    extracted_text = sanitize_text(extracted_text)
                    text_item = {
                        "text": extracted_text,
                        "is_title": shape == slide.shapes.title if hasattr(slide.shapes, 'title') else False
                    }
                    slide_content["texts"].append(text_item)

                    # Check for undefined abbreviations
                    words = extracted_text.split()
                    for word in words:
                        if word.isupper() and len(word) > 1:  # Simplistic abbreviation check
                            undefined_abbreviations.add(word)

                # Note if shape has image
                if hasattr(shape, "image"):
                    slide_content["shapes"].append({
                        "type": "image",
                        "name": shape.name
                    })

            content["slides"].append(slide_content)

        # Prompt user for undefined abbreviations
        self.prompt_for_abbreviation_definitions(undefined_abbreviations, 'data/abbreviations.json', content)

        return content
    
    def save_as_json(self, output_path):
        """Save extracted content as JSON for inspection"""
        content = self.extract_all_content()
        output_path = Path(output_path)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(content, f, indent=2, ensure_ascii=False)
        
        print(f"Content saved to {output_path}")
        return content